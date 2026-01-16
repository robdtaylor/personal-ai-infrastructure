#!/usr/bin/env python3
"""
MS Office Document Generator

Generates Word (.docx) and PowerPoint (.pptx) documents from structured content.

Usage:
    python generate.py --type word --title "Report Title" --content '{"sections": [...]}' --output ~/Downloads/report.docx
    python generate.py --type powerpoint --title "Presentation" --content '{"slides": [...]}' --output ~/Downloads/slides.pptx
"""

import argparse
import json
import os
import sys
from pathlib import Path

# Word document support
from docx import Document
from docx.shared import Inches, Pt
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.enum.style import WD_STYLE_TYPE

# PowerPoint support
from pptx import Presentation
from pptx.util import Inches as PptxInches, Pt as PptxPt
from pptx.enum.text import PP_ALIGN


def create_word_document(title: str, content: dict, output_path: str) -> str:
    """
    Create a Word document from structured content.

    Content format:
    {
        "subtitle": "Optional subtitle",
        "author": "Author name",
        "sections": [
            {
                "heading": "Section Title",
                "level": 1,  # 1-3 for heading levels
                "content": "Paragraph text or list of paragraphs",
                "bullets": ["item1", "item2"],  # Optional bullet list
                "table": {  # Optional table
                    "headers": ["Col1", "Col2"],
                    "rows": [["val1", "val2"], ["val3", "val4"]]
                }
            }
        ]
    }
    """
    doc = Document()

    # Set document title
    title_para = doc.add_heading(title, level=0)
    title_para.alignment = WD_ALIGN_PARAGRAPH.CENTER

    # Add subtitle if provided
    if content.get("subtitle"):
        subtitle = doc.add_paragraph(content["subtitle"])
        subtitle.alignment = WD_ALIGN_PARAGRAPH.CENTER

    # Add author if provided
    if content.get("author"):
        author = doc.add_paragraph(f"By {content['author']}")
        author.alignment = WD_ALIGN_PARAGRAPH.CENTER
        doc.add_paragraph()  # Spacer

    # Process sections
    for section in content.get("sections", []):
        # Add heading
        if section.get("heading"):
            level = section.get("level", 1)
            doc.add_heading(section["heading"], level=min(level, 3))

        # Add image if provided
        if section.get("image"):
            image_path = os.path.expanduser(section["image"])
            if os.path.exists(image_path):
                # Default width is 6 inches, can be customized
                width = Inches(section.get("image_width", 6))
                doc.add_picture(image_path, width=width)
                # Center the image
                last_paragraph = doc.paragraphs[-1]
                last_paragraph.alignment = WD_ALIGN_PARAGRAPH.CENTER
                # Add caption if provided
                if section.get("image_caption"):
                    caption = doc.add_paragraph(section["image_caption"])
                    caption.alignment = WD_ALIGN_PARAGRAPH.CENTER
                    caption.runs[0].italic = True if caption.runs else None

        # Add paragraph content
        if section.get("content"):
            if isinstance(section["content"], list):
                for para in section["content"]:
                    doc.add_paragraph(para)
            else:
                doc.add_paragraph(section["content"])

        # Add bullet list
        if section.get("bullets"):
            for bullet in section["bullets"]:
                doc.add_paragraph(bullet, style='List Bullet')

        # Add numbered list
        if section.get("numbered"):
            for item in section["numbered"]:
                doc.add_paragraph(item, style='List Number')

        # Add table
        if section.get("table"):
            table_data = section["table"]
            headers = table_data.get("headers", [])
            rows = table_data.get("rows", [])

            if headers or rows:
                num_cols = len(headers) if headers else len(rows[0]) if rows else 0
                num_rows = (1 if headers else 0) + len(rows)

                if num_cols > 0 and num_rows > 0:
                    table = doc.add_table(rows=num_rows, cols=num_cols)
                    table.style = 'Table Grid'

                    # Add headers
                    if headers:
                        header_row = table.rows[0]
                        for i, header in enumerate(headers):
                            cell = header_row.cells[i]
                            cell.text = str(header)
                            # Bold headers
                            for para in cell.paragraphs:
                                for run in para.runs:
                                    run.bold = True

                    # Add data rows
                    start_row = 1 if headers else 0
                    for i, row_data in enumerate(rows):
                        row = table.rows[start_row + i]
                        for j, cell_data in enumerate(row_data):
                            if j < num_cols:
                                row.cells[j].text = str(cell_data)

                doc.add_paragraph()  # Spacer after table

    # Save document
    doc.save(output_path)
    return output_path


def create_powerpoint(title: str, content: dict, output_path: str) -> str:
    """
    Create a PowerPoint presentation from structured content.

    Content format:
    {
        "subtitle": "Optional subtitle",
        "author": "Author name",
        "slides": [
            {
                "title": "Slide Title",
                "layout": "title|content|two_column|section|blank",
                "content": "Text content or list of bullet points",
                "bullets": ["point1", "point2"],
                "left_column": ["items"],  # For two_column layout
                "right_column": ["items"],  # For two_column layout
                "notes": "Speaker notes"
            }
        ]
    }
    """
    prs = Presentation()
    prs.slide_width = PptxInches(13.333)  # 16:9 aspect ratio
    prs.slide_height = PptxInches(7.5)

    # Layout indices (standard PowerPoint layouts)
    LAYOUT_TITLE = 0
    LAYOUT_TITLE_CONTENT = 1
    LAYOUT_SECTION = 2
    LAYOUT_TWO_CONTENT = 3
    LAYOUT_BLANK = 6

    # Add title slide
    title_slide_layout = prs.slide_layouts[LAYOUT_TITLE]
    slide = prs.slides.add_slide(title_slide_layout)
    slide.shapes.title.text = title
    if slide.placeholders[1]:
        subtitle_text = content.get("subtitle", "")
        if content.get("author"):
            subtitle_text = f"{subtitle_text}\n{content['author']}" if subtitle_text else content["author"]
        slide.placeholders[1].text = subtitle_text

    # Process slides
    for slide_data in content.get("slides", []):
        layout_name = slide_data.get("layout", "content")

        # Select layout
        if layout_name == "title":
            layout = prs.slide_layouts[LAYOUT_TITLE]
        elif layout_name == "section":
            layout = prs.slide_layouts[LAYOUT_SECTION]
        elif layout_name == "two_column":
            layout = prs.slide_layouts[LAYOUT_TWO_CONTENT]
        elif layout_name == "blank":
            layout = prs.slide_layouts[LAYOUT_BLANK]
        else:  # Default to content layout
            layout = prs.slide_layouts[LAYOUT_TITLE_CONTENT]

        slide = prs.slides.add_slide(layout)

        # Set title
        if slide.shapes.title and slide_data.get("title"):
            slide.shapes.title.text = slide_data["title"]

        # Handle different layouts
        if layout_name == "two_column":
            # Two column content
            left_items = slide_data.get("left_column", [])
            right_items = slide_data.get("right_column", [])

            # Find content placeholders
            placeholders = [p for p in slide.placeholders if p.placeholder_format.idx > 0]
            if len(placeholders) >= 2:
                # Left column
                if left_items:
                    tf = placeholders[0].text_frame
                    tf.clear()
                    for i, item in enumerate(left_items):
                        if i == 0:
                            tf.paragraphs[0].text = item
                        else:
                            p = tf.add_paragraph()
                            p.text = item
                        tf.paragraphs[i].level = 0

                # Right column
                if right_items:
                    tf = placeholders[1].text_frame
                    tf.clear()
                    for i, item in enumerate(right_items):
                        if i == 0:
                            tf.paragraphs[0].text = item
                        else:
                            p = tf.add_paragraph()
                            p.text = item
                        tf.paragraphs[i].level = 0

        elif layout_name not in ["title", "section", "blank"]:
            # Standard content slide
            content_placeholder = None
            for placeholder in slide.placeholders:
                if placeholder.placeholder_format.idx == 1:
                    content_placeholder = placeholder
                    break

            if content_placeholder:
                tf = content_placeholder.text_frame
                tf.clear()

                # Get content (bullets or content field)
                items = slide_data.get("bullets", [])
                if not items and slide_data.get("content"):
                    if isinstance(slide_data["content"], list):
                        items = slide_data["content"]
                    else:
                        items = [slide_data["content"]]

                for i, item in enumerate(items):
                    if i == 0:
                        tf.paragraphs[0].text = item
                    else:
                        p = tf.add_paragraph()
                        p.text = item
                    tf.paragraphs[i].level = 0

        # Add image to slide if provided
        if slide_data.get("image"):
            image_path = os.path.expanduser(slide_data["image"])
            if os.path.exists(image_path):
                # Position image - default centered below title
                left = PptxInches(slide_data.get("image_left", 1.5))
                top = PptxInches(slide_data.get("image_top", 2))
                width = PptxInches(slide_data.get("image_width", 10))
                slide.shapes.add_picture(image_path, left, top, width=width)

        # Add speaker notes
        if slide_data.get("notes"):
            notes_slide = slide.notes_slide
            notes_slide.notes_text_frame.text = slide_data["notes"]

    # Save presentation
    prs.save(output_path)
    return output_path


def main():
    parser = argparse.ArgumentParser(description="Generate MS Office documents")
    parser.add_argument("--type", "-t", required=True, choices=["word", "powerpoint"],
                       help="Document type to generate")
    parser.add_argument("--title", required=True, help="Document title")
    parser.add_argument("--content", required=True,
                       help="JSON content structure or path to JSON file")
    parser.add_argument("--output", "-o", required=True, help="Output file path")

    args = parser.parse_args()

    # Parse content
    if os.path.isfile(args.content):
        with open(args.content, 'r') as f:
            content = json.load(f)
    else:
        try:
            content = json.loads(args.content)
        except json.JSONDecodeError as e:
            print(f"Error parsing content JSON: {e}", file=sys.stderr)
            sys.exit(1)

    # Expand output path
    output_path = os.path.expanduser(args.output)

    # Ensure output directory exists
    os.makedirs(os.path.dirname(output_path) or '.', exist_ok=True)

    # Generate document
    try:
        if args.type == "word":
            result = create_word_document(args.title, content, output_path)
            print(f"Word document created: {result}")
        else:
            result = create_powerpoint(args.title, content, output_path)
            print(f"PowerPoint presentation created: {result}")
    except Exception as e:
        print(f"Error creating document: {e}", file=sys.stderr)
        sys.exit(1)


if __name__ == "__main__":
    main()
